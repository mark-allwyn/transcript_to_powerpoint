"""
Transcript → Slide Deck
-----------------------
AI-powered application that converts meeting transcripts into PowerPoint presentations
with AI-generated images using OpenAI's structured output and DALL-E 3.

Run: streamlit run app.py
"""

import os, json, math, requests, base64, time
from io import BytesIO
from typing import List
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field
import openai

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
openai.api_key = os.getenv("OPENAI_API_KEY")

# ---------------------------------------------------------------------------
# Pydantic Models for Structured Output
# ---------------------------------------------------------------------------
class MeetingSummary(BaseModel):
    """Structured output for meeting transcript summary"""
    key_points: List[str] = Field(description="List of 3-7 important discussion points from the meeting")
    decisions: List[str] = Field(description="List of decisions made during the meeting", default=[])
    action_items: List[str] = Field(description="List of action items identified", default=[])

class SlideSpec(BaseModel):
    """Specification for a single slide"""
    title: str = Field(description="Slide title, max 8 words, present tense")
    bullets: List[str] = Field(description="3-6 concise bullet points, each under 80 characters")

class SlideSpecs(BaseModel):
    """Collection of slide specifications"""
    slides: List[SlideSpec] = Field(description="Array of slide specifications")

class ImagePrompts(BaseModel):
    """Collection of image generation prompts"""
    prompts: List[str] = Field(description="Array of DALL-E prompts, one for each slide")

# ---------------------------------------------------------------------------
# Helper Functions
# ---------------------------------------------------------------------------
def chunk_text(text: str, words_per_chunk: int = 8_000) -> list[str]:
    """Split very long transcripts into chunks for processing."""
    words = text.split()
    if len(words) <= words_per_chunk:
        return [text]
    n_chunks = math.ceil(len(words) / words_per_chunk)
    return [
        " ".join(words[i::n_chunks])
        for i in range(n_chunks)
    ]


def merge_summaries(parts: list[dict]) -> dict:
    merged = {"key_points": [], "decisions": [], "action_items": []}
    for p in parts:
        for k in merged:
            merged[k].extend(p.get(k, []))
    return merged


def create_placeholder_image() -> bytes:
    """Create a simple white placeholder image."""
    try:
        from PIL import Image
        img = Image.new('RGB', (1024, 1024), 'white')
        buf = BytesIO()
        img.save(buf, format='PNG')
        return buf.getvalue()
    except ImportError:
        # Fallback: small white PNG image
        white_image_b64 = """
        iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAGklEQVQYV2P4//8/AzYwOjraMDo6
        2jA6OtoAALvMBONcfOGOAAAAAElFTkSuQmCC
        """
        return base64.b64decode(white_image_b64.replace('\n', '').replace(' ', ''))


def create_images_gpt(prompts: list[str]) -> list[bytes]:
    """Generate images using DALL-E 3 and return raw bytes."""
    bins = []
    for prompt in prompts:
        try:
            resp = openai.images.generate(
                model="dall-e-3",
                prompt=prompt,
                n=1,
                size="1024x1024",
                quality="standard",
            )
            url = resp.data[0].url
            if url:
                image_data = requests.get(url, timeout=30).content
                bins.append(image_data)
            else:
                bins.append(create_placeholder_image())
        except Exception:
            bins.append(create_placeholder_image())
    return bins


def build_pptx(slide_specs: list[dict], images: list[bytes]) -> BytesIO:
    """Build PowerPoint presentation from slide specifications and images."""
    prs = Presentation()
    
    # Title slide
    tslide = prs.slides.add_slide(prs.slide_layouts[0])
    tslide.shapes.title.text = "Meeting Summary"
    tslide.placeholders[1].text = "AI-generated deck"

    # Content slides
    bullet_layout = prs.slide_layouts[1]
    for spec, img_bytes in zip(slide_specs, images):
        sld = prs.slides.add_slide(bullet_layout)
        sld.shapes.title.text = spec["title"]

        body = sld.shapes.placeholders[1].text_frame
        body.clear()
        for bullet in spec["bullets"]:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

        # Add image on right side
        pic_stream = BytesIO(img_bytes)
        sld.shapes.add_picture(pic_stream, Inches(5.5), Inches(1.3), width=Inches(3))

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Main Processing Pipeline
# ---------------------------------------------------------------------------

def generate_slide_package(transcript: str):
    """Process transcript and generate slide specifications with images."""
    from openai import OpenAI
    client = OpenAI()
    
    start_time = time.time()
    
    # Step 1: Analyze transcript using OpenAI structured output
    step1_start = time.time()
    try:
        summary_response = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system", 
                    "content": "You are an expert meeting transcript analyzer. Extract key information from meeting transcripts accurately and completely."
                },
                {
                    "role": "user",
                    "content": f"""
                    Analyze this meeting transcript and extract key information.
                    
                    Extract only what was actually discussed in the meeting:
                    - Key discussion points (3-7 most important items)
                    - Decisions made during the meeting (actual decisions)
                    - Action items identified (specific tasks mentioned)
                    
                    Meeting transcript:
                    {transcript}
                    """
                }
            ],
            response_format=MeetingSummary
        )
        summary_json = summary_response.choices[0].message.parsed.model_dump()
    except Exception:
        # Fallback to manual parsing
        summary_json = {"key_points": ["Meeting analysis failed"], "decisions": [], "action_items": []}
    
    step1_time = time.time() - step1_start

    # Step 2: Create slide specifications
    step2_start = time.time()
    try:
        slides_response = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "You are a presentation expert who creates slide specifications based on meeting content."
                },
                {
                    "role": "user",
                    "content": f"""
                    Create slide specifications based on this meeting summary.
                    
                    MEETING SUMMARY:
                    {json.dumps(summary_json, indent=2)}
                    
                    Requirements:
                    - Create slides specific to this meeting content
                    - Titles max 8 words, present tense
                    - 3-6 bullet points per slide, each under 80 characters
                    - Use the actual meeting information provided
                    - Professional business presentation style
                    
                    Create slides covering key points, decisions, and action items from the meeting.
                    """
                }
            ],
            response_format=SlideSpecs
        )
        slide_specs = slides_response.choices[0].message.parsed.slides
        slide_specs_data = [spec.model_dump() for spec in slide_specs]
    except Exception:
        # Fallback: Create slides directly from summary
        slide_specs_data = []
        if summary_json.get('key_points'):
            slide_specs_data.append({
                "title": "Key Discussion Points",
                "bullets": summary_json['key_points'][:6]
            })
        if summary_json.get('decisions'):
            slide_specs_data.append({
                "title": "Decisions Made", 
                "bullets": summary_json['decisions'][:6]
            })
        if summary_json.get('action_items'):
            slide_specs_data.append({
                "title": "Action Items",
                "bullets": summary_json['action_items'][:6]
            })
        
        if not slide_specs_data:
            slide_specs_data = [{"title": "Meeting Summary", "bullets": ["Content not available"]}]
    
    step2_time = time.time() - step2_start

    # Step 3: Generate image prompts
    step3_start = time.time()
    try:
        prompts_response = client.beta.chat.completions.parse(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are a visual design expert who creates DALL-E prompts for business presentations."
                },
                {
                    "role": "user",
                    "content": f"""
                    Create DALL-E image prompts for these slides based on the meeting content.
                    
                    MEETING CONTEXT:
                    {json.dumps(summary_json, indent=2)}
                    
                    SLIDES TO CREATE IMAGES FOR:
                    {json.dumps(slide_specs_data, indent=2)}
                    
                    Requirements:
                    - Create professional business illustrations
                    - Modern vector-style graphics in blue/gray/white colors
                    - Relevant to the specific meeting content and slide topics
                    - ABSOLUTELY NO TEXT, WORDS, LETTERS, OR NUMBERS in images
                    - Minimalist, clean design
                    - End each prompt with ", no text, no words, no labels"
                    
                    Create one detailed prompt per slide that relates to the actual meeting topics.
                    """
                }
            ],
            response_format=ImagePrompts
        )
        prompts = prompts_response.choices[0].message.parsed.prompts
    except Exception:
        prompts = [f"Minimalist business illustration for slide {i+1}, no text, no words, no labels" for i in range(len(slide_specs_data))]
    
    step3_time = time.time() - step3_start

    # Ensure prompts match slides count
    if len(prompts) != len(slide_specs_data):
        if len(prompts) < len(slide_specs_data):
            prompts.extend([f"Business illustration, no text" for _ in range(len(slide_specs_data) - len(prompts))])
        else:
            prompts = prompts[:len(slide_specs_data)]

    # Step 4: Generate images
    step4_start = time.time()
    image_bins = create_images_gpt(prompts)
    step4_time = time.time() - step4_start
    
    total_time = time.time() - start_time

    return slide_specs_data, image_bins, {
        "transcript_analysis": step1_time,
        "slide_generation": step2_time, 
        "image_prompts": step3_time,
        "image_generation": step4_time,
        "total_time": total_time
    }


# ---------------------------------------------------------------------------
# Streamlit User Interface
# ---------------------------------------------------------------------------
st.title("Transcript → Slides")
file = st.file_uploader("Upload meeting transcript (.txt)", type=["txt"])

if file:
    transcript_text = file.read().decode("utf-8", errors="ignore")
    with st.spinner("Processing transcript…"):
        specs, imgs, timing_info = generate_slide_package(transcript_text)
        deck = build_pptx(specs, imgs)
    
    st.success("Slide deck ready!")
    
    # Display timing information
    st.subheader("Processing Times")
    col1, col2 = st.columns(2)
    
    with col1:
        st.metric("Transcript Analysis", f"{timing_info['transcript_analysis']:.1f}s")
        st.metric("Slide Generation", f"{timing_info['slide_generation']:.1f}s")
    
    with col2:
        st.metric("Image Prompts", f"{timing_info['image_prompts']:.1f}s")
        st.metric("Image Generation", f"{timing_info['image_generation']:.1f}s")
    
    st.metric("Total Processing Time", f"{timing_info['total_time']:.1f}s")
    
    st.download_button(
        label="Download PPTX",
        data=deck,
        file_name="meeting_summary.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
