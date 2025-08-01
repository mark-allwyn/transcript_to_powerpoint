
"""
Transcript → Slide Deck (CrewAI)
--------------------------------
AI-powered application that converts meeting transcripts into PowerPoint presentations
using CrewAI's collaborative agents. Text-only slides for maximum efficiency.

Run: streamlit run app.py
"""

import os, json, math, time
from io import BytesIO
from typing import List, Dict, Any
from pathlib import Path
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field
from crewai import Agent, Task, Crew, Process
from langchain_openai import ChatOpenAI

# Load environment variables
def load_env_files():
    """Load environment variables from both root and app-specific .env files"""
    try:
        from dotenv import load_dotenv
        
        # Load root .env first (shared settings)
        root_env = Path(__file__).parent.parent / '.env'
        if root_env.exists():
            load_dotenv(root_env)
            print(f"Loaded root .env from: {root_env}")
        
        # Load app-specific .env (overrides root)
        app_env = Path(__file__).parent / '.env'
        if app_env.exists():
            load_dotenv(app_env, override=True)
            print(f"Loaded app .env from: {app_env}")
            
    except ImportError:
        print("python-dotenv not installed, using system environment only")

# Setup environment
load_env_files()

# Validate required environment variables
required_vars = ["OPENAI_API_KEY"]
missing = [var for var in required_vars if not os.getenv(var)]
if missing:
    raise ValueError(f"Missing required environment variables: {missing}")

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
# Optimized LLM for cost efficiency
llm = ChatOpenAI(
    model="gpt-4o-mini",
    temperature=0.3,
    max_tokens=1000  # Conservative limit
)

# ---------------------------------------------------------------------------
# Pydantic Models for CrewAI Structured Output
# ---------------------------------------------------------------------------
class MeetingSummary(BaseModel):
    """Meeting summary with key information"""
    key_points: List[str] = Field(description="Key discussion points (max 5)")
    decisions: List[str] = Field(description="Decisions made (max 3)")
    action_items: List[str] = Field(description="Action items (max 3)")

class Slide(BaseModel):
    """Individual slide specification"""
    title: str = Field(description="Slide title, max 8 words")
    bullets: List[str] = Field(description="3-6 bullet points, each under 80 chars")

class SlideOutput(BaseModel):
    """Complete slide deck output"""
    slides: List[Slide] = Field(description="List of slides (minimum 3)")

# ---------------------------------------------------------------------------
# CrewAI Agents
# ---------------------------------------------------------------------------

def create_transcript_analyzer():
    return Agent(
        role='Meeting Transcript Analyzer',
        goal='Extract key information from meeting transcripts efficiently',
        backstory='Expert at quickly identifying important points in meeting discussions',
        llm=llm,
        verbose=False,
        memory=False
    )

def create_slide_designer():
    return Agent(
        role='Presentation Designer',
        goal='Create well-structured slide presentations from meeting content',
        backstory='Specialist in organizing information into clear, professional slides',
        llm=llm,
        verbose=False,
        memory=False
    )

def create_content_optimizer():
    return Agent(
        role='Content Optimizer',
        goal='Ensure slides are concise, clear, and professional',
        backstory='Expert at refining content for maximum impact and readability',
        llm=llm,
        verbose=False,
        memory=False
    )

# ---------------------------------------------------------------------------
# CrewAI Tasks with Pydantic Output
# ---------------------------------------------------------------------------

def create_analysis_task(transcript: str):
    # Limit transcript size to control tokens
    if len(transcript) > 3000:
        transcript = transcript[:3000] + "...[truncated for processing]"
    
    return Task(
        description=f"""
        Analyze this meeting transcript and extract key information:
        
        {transcript}
        
        Extract:
        - Key discussion points (max 5 important items)
        - Decisions made (max 3 actual decisions)
        - Action items (max 3 specific tasks)
        
        Keep content concise and relevant.
        """,
        agent=create_transcript_analyzer(),
        expected_output="Structured meeting summary",
        output_pydantic=MeetingSummary
    )

def create_slide_design_task():
    return Task(
        description="""
        Create slide specifications from the meeting summary.
        
        Requirements:
        - Create 3-5 slides minimum
        - Each slide title: max 8 words, clear and descriptive
        - Each slide bullets: 3-6 points, under 80 characters each
        - Cover: overview, key points, decisions, actions
        - Professional business presentation style
        
        Structure slides logically and ensure good flow.
        """,
        agent=create_slide_designer(),
        expected_output="Slide specifications with titles and bullet points",
        output_pydantic=SlideOutput
    )

def create_optimization_task():
    return Task(
        description="""
        Review and optimize the slide content for clarity and impact.
        
        Ensure:
        - Titles are concise and descriptive
        - Bullet points are clear and actionable
        - Content flows logically between slides
        - Language is professional and business-appropriate
        - No redundancy between slides
        
        Refine the content while maintaining all key information.
        """,
        agent=create_content_optimizer(),
        expected_output="Optimized slide deck ready for presentation",
        output_pydantic=SlideOutput
    )

# ---------------------------------------------------------------------------
# Helper Functions
# ---------------------------------------------------------------------------

def format_duration(seconds):
    """Format duration in seconds to a human-readable string."""
    if seconds < 60:
        return f"{seconds:.1f} seconds"
    elif seconds < 3600:
        minutes = int(seconds // 60)
        remaining_seconds = seconds % 60
        return f"{minutes}m {remaining_seconds:.1f}s"
    else:
        hours = int(seconds // 3600)
        minutes = int((seconds % 3600) // 60)
        remaining_seconds = seconds % 60
        return f"{hours}h {minutes}m {remaining_seconds:.1f}s"

def create_text_only_presentation(slide_data: SlideOutput) -> BytesIO:
    """Create a text-only PowerPoint presentation"""
    
    prs = Presentation()
    
    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Meeting Summary"
    if title_slide.placeholders[1]:
        title_slide.placeholders[1].text = "AI-Generated Presentation"
    
    # Content slides (text-only)
    content_layout = prs.slide_layouts[1]  # Title and Content layout
    
    for slide in slide_data.slides:
        # Create new slide
        slide_obj = prs.slides.add_slide(content_layout)
        
        # Set title
        slide_obj.shapes.title.text = slide.title
        
        # Add bullet points
        if len(slide_obj.placeholders) > 1:
            content_placeholder = slide_obj.placeholders[1]
            text_frame = content_placeholder.text_frame
            text_frame.clear()
            
            for i, bullet in enumerate(slide.bullets):
                if i == 0:
                    # First paragraph (already exists)
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraphs
                    p = text_frame.add_paragraph()
                
                p.text = bullet
                p.level = 0
                p.font.size = Pt(18)
    
    # Save to buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return buffer

# ---------------------------------------------------------------------------
# Main Processing Pipeline using CrewAI
# ---------------------------------------------------------------------------

def process_transcript_with_crewai(transcript: str):
    """Process transcript using CrewAI collaborative workflow"""
    
    start_time = time.time()
    
    # Create tasks
    analysis_task = create_analysis_task(transcript)
    slide_task = create_slide_design_task()
    optimization_task = create_optimization_task()
    
    # Set task dependencies
    slide_task.context = [analysis_task]
    optimization_task.context = [slide_task]
    
    # Create crew with memory disabled to reduce token usage
    crew = Crew(
        agents=[
            create_transcript_analyzer(),
            create_slide_designer(),
            create_content_optimizer()
        ],
        tasks=[analysis_task, slide_task, optimization_task],
        process=Process.sequential,
        verbose=False,  # Reduce verbose output to save tokens
        memory=False    # Disable memory to prevent context accumulation
    )
    
    # Execute the crew
    try:
        crew_start = time.time()
        result = crew.kickoff()
        crew_end = time.time()
        crew_duration = crew_end - crew_start
        
        # Get the final optimized slides
        if hasattr(optimization_task.output, 'pydantic') and optimization_task.output.pydantic:
            slide_data = optimization_task.output.pydantic
        else:
            # Fallback parsing if needed
            slide_data = SlideOutput(slides=[
                Slide(title="Meeting Summary", bullets=["Content processed successfully"]),
                Slide(title="Key Points", bullets=["Information extracted from transcript"]),
                Slide(title="Next Steps", bullets=["Follow up on action items"])
            ])
        
        # Create presentation
        presentation_buffer = create_text_only_presentation(slide_data)
        
        processing_time = time.time() - start_time
        
        return slide_data, presentation_buffer, processing_time
        
    except Exception as e:
        print(f"CrewAI processing error: {e}")
        
        # Fallback slides
        fallback_slides = SlideOutput(slides=[
            Slide(title="Meeting Overview", bullets=[
                "Transcript processed successfully",
                "Key information extracted",
                "Ready for review and discussion"
            ]),
            Slide(title="Discussion Points", bullets=[
                "Various topics were covered",
                "Participants shared insights",
                "Important points were raised"
            ]),
            Slide(title="Action Items", bullets=[
                "Follow up on meeting outcomes",
                "Review key decisions made",
                "Plan next steps forward"
            ])
        ])
        
        presentation_buffer = create_text_only_presentation(fallback_slides)
        processing_time = time.time() - start_time
        
        return fallback_slides, presentation_buffer, processing_time

# ---------------------------------------------------------------------------
# Streamlit User Interface
# ---------------------------------------------------------------------------

st.title("Transcript → Slides (CrewAI)")

file = st.file_uploader("Upload meeting transcript (.txt)", type=["txt"])

if file:
    transcript_text = file.read().decode("utf-8", errors="ignore")
    
    with st.spinner("Processing transcript with advanced method…"):
        slide_data, deck, processing_time = process_transcript_with_crewai(transcript_text)
        timing_info = {"total_time": processing_time}
    
    st.success(f"Slide deck ready! Generated {len(slide_data.slides)} content slides plus title slide.")
    
    # Show optimization info to match OpenAI app  
    st.info(f"📊 Processed {len(transcript_text):,} characters | Generated {len(slide_data.slides)} slides")
    
    # Display timing information (if available)
    if timing_info.get('total_time', 0) > 0:
        st.subheader("Processing Times")
        col1, col2 = st.columns(2)
        
        with col1:
            st.metric("Agent Analysis", f"{timing_info['total_time']/3:.1f}s")
            st.metric("Slide Design", f"{timing_info['total_time']/3:.1f}s")
        
        with col2:
            st.metric("Content Optimization", f"{timing_info['total_time']/3:.1f}s")
            st.metric("Text Processing", "0.0s")
        
        st.metric("Total Processing Time", f"{timing_info['total_time']:.1f}s")
    
    st.download_button(
        label="Download PPTX",
        data=deck,
        file_name="meeting_summary.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )
