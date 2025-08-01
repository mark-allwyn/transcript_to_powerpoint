"""
Transcript → Slide Deck
-----------------------
AI-powered application that converts meeting transcripts into PowerPoint presentations
with AI-generated images using OpenAI's structured output and DALL-E 3.

Run: streamlit run app.py
"""

import os, json, math, requests, base64, time
from io import BytesIO
from typing import List
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pydantic import BaseModel, Field
import openai

# Load environment variables using shared utility
import sys
from pathlib import Path

# Load environment variables
def load_env_files():
    """Load environment variables from both root and app-specific .env files"""
    try:
        from dotenv import load_dotenv
        
        # Load root .env first (shared settings)
        root_env = Path(__file__).parent.parent / '.env'
        if root_env.exists():
            load_dotenv(root_env)
            print(f"Loaded root .env from: {root_env}")
        
        # Load app-specific .env (overrides root)
        app_env = Path(__file__).parent / '.env'
        if app_env.exists():
            load_dotenv(app_env, override=True)
            print(f"Loaded app .env from: {app_env}")
            
    except ImportError:
        print("python-dotenv not installed, using system environment only")

# Setup environment
load_env_files()

# Validate required environment variables
required_vars = ["OPENAI_API_KEY"]
missing = [var for var in required_vars if not os.getenv(var)]
if missing:
    raise ValueError(f"Missing required environment variables: {missing}")

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------
openai.api_key = os.getenv("OPENAI_API_KEY")

# ---------------------------------------------------------------------------
# Pydantic Models for Structured Output
# ---------------------------------------------------------------------------
class MeetingSummary(BaseModel):
    """Structured output for meeting transcript summary"""
    key_points: List[str] = Field(description="List of 3-7 important discussion points from the meeting")
    decisions: List[str] = Field(description="List of decisions made during the meeting", default=[])
    action_items: List[str] = Field(description="List of action items identified", default=[])

class SlideSpec(BaseModel):
    """Specification for a single slide"""
    title: str = Field(description="Slide title, max 8 words, present tense")
    bullets: List[str] = Field(description="3-6 concise bullet points, each under 80 characters")

class SlideSpecs(BaseModel):
    """Collection of slide specifications"""
    slides: List[SlideSpec] = Field(description="Array of slide specifications")

class ImagePrompts(BaseModel):
    """Collection of image generation prompts"""
    prompts: List[str] = Field(description="Array of DALL-E prompts, one for each slide")

# ---------------------------------------------------------------------------
# Helper Functions
# ---------------------------------------------------------------------------
def chunk_text(text: str, words_per_chunk: int = 8_000) -> list[str]:
    """Split very long transcripts into chunks for processing."""
    words = text.split()
    if len(words) <= words_per_chunk:
        return [text]
    n_chunks = math.ceil(len(words) / words_per_chunk)
    return [
        " ".join(words[i::n_chunks])
        for i in range(n_chunks)
    ]


def merge_summaries(parts: list[dict]) -> dict:
    merged = {"key_points": [], "decisions": [], "action_items": []}
    for p in parts:
        for k in merged:
            merged[k].extend(p.get(k, []))
    return merged


def create_placeholder_image() -> bytes:
    """Create a simple white placeholder image."""
    try:
        from PIL import Image
        img = Image.new('RGB', (1024, 1024), 'white')
        buf = BytesIO()
        img.save(buf, format='PNG')
        return buf.getvalue()
    except ImportError:
        # Fallback: small white PNG image
        white_image_b64 = """
        iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAAGklEQVQYV2P4//8/AzYwOjraMDo6
        2jA6OtoAALvMBONcfOGOAAAAAElFTkSuQmCC
        """
        return base64.b64decode(white_image_b64.replace('\n', '').replace(' ', ''))


def create_images_gpt(prompts: list[str]) -> list[bytes]:
    """Generate images using DALL-E 3 and return raw bytes."""
    bins = []
    for prompt in prompts:
        try:
            resp = openai.images.generate(
                model="dall-e-3",
                prompt=prompt,
                n=1,
                size="1024x1024",
                quality="standard",
            )
            url = resp.data[0].url
            if url:
                image_data = requests.get(url, timeout=30).content
                bins.append(image_data)
            else:
                bins.append(create_placeholder_image())
        except Exception:
            bins.append(create_placeholder_image())
    return bins


def build_pptx(slide_specs: list[dict], images: list[bytes]) -> BytesIO:
    """Build PowerPoint presentation from slide specifications and images."""
    prs = Presentation()
    
    print(f"Building PowerPoint with {len(slide_specs)} slides")
    
    # Title slide
    tslide = prs.slides.add_slide(prs.slide_layouts[0])
    tslide.shapes.title.text = "Meeting Summary"
    tslide.placeholders[1].text = "AI-generated deck"

    # Content slides
    bullet_layout = prs.slide_layouts[1]
    for i, (spec, img_bytes) in enumerate(zip(slide_specs, images)):
        print(f"Creating slide {i+2}: {spec['title']}")
        sld = prs.slides.add_slide(bullet_layout)
        sld.shapes.title.text = spec["title"]

        body = sld.shapes.placeholders[1].text_frame
        body.clear()
        for bullet in spec["bullets"]:
            p = body.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)

        # Add image on right side
        try:
            pic_stream = BytesIO(img_bytes)
            sld.shapes.add_picture(pic_stream, Inches(5.5), Inches(1.3), width=Inches(3))
        except Exception as e:
            print(f"Error adding image to slide {i+2}: {e}")

    print(f"Final presentation has {len(prs.slides)} total slides")
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Main Processing Pipeline
# ---------------------------------------------------------------------------

def generate_slide_package(transcript: str):
    """Process transcript and generate slide specifications with images."""
    from openai import OpenAI
    client = OpenAI()
    
    start_time = time.time()
    
    # OPTIMIZATION: Limit transcript length to reduce tokens
    max_chars = 8000  # Roughly 2000 tokens
    if len(transcript) > max_chars:
        print(f"Truncating transcript from {len(transcript)} to {max_chars} characters")
        transcript = transcript[:max_chars] + "...\n[Content truncated for processing efficiency]"
    
    # Step 1: Analyze transcript using OpenAI structured output (OPTIMIZED)
    step1_start = time.time()
    try:
        summary_response = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system", 
                    "content": "Extract key meeting information concisely. Limit each category to essential points only."
                },
                {
                    "role": "user",
                    "content": f"""
                    Extract key information from this transcript:
                    
                    - Key points (max 5 important items)
                    - Decisions made (max 3 actual decisions)  
                    - Action items (max 3 specific tasks)
                    
                    Transcript: {transcript}
                    """
                }
            ],
            response_format=MeetingSummary
        )
        summary_json = summary_response.choices[0].message.parsed.model_dump()
    except Exception as e:
        print(f"Step 1 error: {e}")
        # Fallback to manual parsing
        summary_json = {"key_points": ["Meeting analysis failed"], "decisions": [], "action_items": []}
    
    step1_time = time.time() - step1_start

    # Step 2: Create slide specifications with explicit slide structure (OPTIMIZED)
    step2_start = time.time()
    try:
        # Create a condensed summary for slide generation (reduce tokens)
        condensed_summary = {
            "key_points": summary_json.get('key_points', [])[:5],  # Limit to 5 points
            "decisions": summary_json.get('decisions', [])[:3],    # Limit to 3 decisions  
            "action_items": summary_json.get('action_items', [])[:3] # Limit to 3 actions
        }
        
        slides_response = client.beta.chat.completions.parse(
            model="gpt-4o-mini",
            messages=[
                {
                    "role": "system",
                    "content": "Create multiple slide specifications from meeting summaries. Always generate 3-5 slides minimum."
                },
                {
                    "role": "user",
                    "content": f"""
                    Create 4-5 slides from this summary. Return slides array.
                    
                    Summary: {json.dumps(condensed_summary)}
                    
                    Requirements:
                    - 3-5 slides minimum
                    - Titles: max 8 words
                    - Bullets: 3-6 points, under 80 chars each
                    - Cover: overview, key points, decisions, actions
                    """
                }
            ],
            response_format=SlideSpecs
        )
        slide_specs = slides_response.choices[0].message.parsed.slides
        slide_specs_data = [spec.model_dump() for spec in slide_specs]
        
        # Debug: Print slide count
        print(f"Generated {len(slide_specs_data)} slides from structured output")
        
    except Exception as e:
        print(f"Step 2 error: {e}")
        # Enhanced fallback: Create multiple slides directly from summary
        slide_specs_data = []
        
        # Always create a summary slide first
        all_content = []
        if summary_json.get('key_points'):
            all_content.extend(summary_json['key_points'][:3])
        if summary_json.get('decisions'):
            all_content.extend([f"Decision: {d}" for d in summary_json['decisions'][:2]])
        if summary_json.get('action_items'):
            all_content.extend([f"Action: {a}" for a in summary_json['action_items'][:2]])
            
        if all_content:
            slide_specs_data.append({
                "title": "Meeting Overview",
                "bullets": all_content[:6]
            })
        
        # Create individual slides for each category if they have content
        if summary_json.get('key_points') and len(summary_json['key_points']) > 0:
            slide_specs_data.append({
                "title": "Key Discussion Points",
                "bullets": summary_json['key_points'][:6]
            })
            
        if summary_json.get('decisions') and len(summary_json['decisions']) > 0:
            slide_specs_data.append({
                "title": "Decisions Made", 
                "bullets": summary_json['decisions'][:6]
            })
            
        if summary_json.get('action_items') and len(summary_json['action_items']) > 0:
            slide_specs_data.append({
                "title": "Action Items",
                "bullets": summary_json['action_items'][:6]
            })
        
        # Ensure we have at least 2 slides
        if len(slide_specs_data) < 2:
            # Split content across multiple slides if we only have one
            if slide_specs_data and len(slide_specs_data[0]['bullets']) > 3:
                bullets = slide_specs_data[0]['bullets']
                mid = len(bullets) // 2
                slide_specs_data = [
                    {"title": "Meeting Summary - Part 1", "bullets": bullets[:mid]},
                    {"title": "Meeting Summary - Part 2", "bullets": bullets[mid:]}
                ]
            else:
                # Last resort: create basic slides
                slide_specs_data = [
                    {"title": "Meeting Summary", "bullets": ["Content processed from transcript"]},
                    {"title": "Key Points", "bullets": ["Meeting content available in detail"]}
                ]
        
        print(f"Fallback created {len(slide_specs_data)} slides")
    
    step2_time = time.time() - step2_start

    # Step 3: Generate image prompts (OPTIMIZED - only send slide titles)
    step3_start = time.time()
    try:
        # Extract just the slide titles and key themes for context
        slide_titles = [spec["title"] for spec in slide_specs_data]
        key_themes = summary_json.get('key_points', [])[:3]  # Just top 3 themes
        
        prompts_response = client.beta.chat.completions.parse(
            model="gpt-4o-mini",  # Use cheaper model for image prompts
            messages=[
                {
                    "role": "system",
                    "content": "Create DALL-E prompts for business presentation slides. Generate one prompt per slide title provided."
                },
                {
                    "role": "user",
                    "content": f"""
                    Create image prompts for these slide titles: {', '.join(slide_titles)}
                    
                    Meeting themes: {', '.join(key_themes) if key_themes else 'Business meeting'}
                    
                    Requirements for each prompt:
                    - Professional business illustration
                    - Modern minimalist style
                    - Blue/gray/white color scheme
                    - NO TEXT, WORDS, OR LABELS in images
                    - Each prompt should end with ", no text, no words, no labels"
                    
                    Return {len(slide_titles)} prompts.
                    """
                }
            ],
            response_format=ImagePrompts
        )
        prompts = prompts_response.choices[0].message.parsed.prompts
    except Exception:
        prompts = [f"Minimalist business illustration for slide {i+1}, no text, no words, no labels" for i in range(len(slide_specs_data))]
    
    step3_time = time.time() - step3_start

    # Ensure prompts match slides count
    if len(prompts) != len(slide_specs_data):
        if len(prompts) < len(slide_specs_data):
            prompts.extend([f"Business illustration, no text" for _ in range(len(slide_specs_data) - len(prompts))])
        else:
            prompts = prompts[:len(slide_specs_data)]

    # Step 4: Generate images
    step4_start = time.time()
    image_bins = create_images_gpt(prompts)
    step4_time = time.time() - step4_start
    
    total_time = time.time() - start_time

    return slide_specs_data, image_bins, {
        "transcript_analysis": step1_time,
        "slide_generation": step2_time, 
        "image_prompts": step3_time,
        "image_generation": step4_time,
        "total_time": total_time,
        "transcript_length": len(transcript),
        "slides_generated": len(slide_specs_data)
    }


# ---------------------------------------------------------------------------
# Streamlit User Interface
# ---------------------------------------------------------------------------


if __name__ == "__main__":
    # Add option to use simple generation
    st.title("Transcript to Slides (OpenAI)")
    st.markdown("*Advanced structured output with DALL-E 3 visuals*")

    file = st.file_uploader("Upload meeting transcript (.txt)", type=["txt"], key="standalone_openai_file_uploader")

    if file:
        transcript_text = file.read().decode("utf-8", errors="ignore")
        
        with st.spinner("Processing transcript with advanced structured output..."):
            specs, imgs, timing_info = generate_slide_package(transcript_text)
            deck = build_pptx(specs, imgs)
        
        st.success(f"Slide deck ready! Generated {len(specs)} content slides plus title slide.")
        
        # Show optimization info
        if timing_info.get('transcript_length'):
            st.info(f"Processed {timing_info['transcript_length']:,} characters | Generated {timing_info.get('slides_generated', 0)} slides")
        
        # Display slide preview
        st.subheader("Slide Preview")
        for i, spec in enumerate(specs, 1):
            with st.expander(f"Slide {i + 1}: {spec['title']}"):
                for bullet in spec['bullets']:
                    st.write(f"• {bullet}")
        
        # Display timing information (if available)
        if timing_info.get('total_time', 0) > 0:
            st.subheader("Processing Times")
            col1, col2 = st.columns(2)
            
            with col1:
                st.metric("Transcript Analysis", f"{timing_info.get('transcript_analysis', 0):.1f}s")
                st.metric("Slide Generation", f"{timing_info.get('slide_generation', 0):.1f}s")
            
            with col2:
                st.metric("Image Prompts", f"{timing_info.get('image_prompts', 0):.1f}s")
                st.metric("Image Generation", f"{timing_info.get('image_generation', 0):.1f}s")
            
            st.metric("Total Processing Time", f"{timing_info['total_time']:.1f}s")
        
        st.download_button(
            label="Download PPTX",
            data=deck,
            file_name="meeting_summary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key="standalone_openai_download_button"
        )
